# -----------------------------------------------------------#
# Name: Chunli Tang                                          #
# School: Auburn University                                  #
# -----------------------------------------------------------#

import os
from tkinter import *
from matplotlib import pyplot
import matplotlib.patches as patches
import numpy as np
import pandas as pd
from tkinter import filedialog
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.shapes.autoshape import Shape as shape
from lmfit import Model
import lmfit


class SHG_Processing():
    def __init__(self):
        super().__init__()
        self.run()

    def run(self):

        SHGpath = 'SHG'  # Processed data directory
        isExist = os.path.exists(SHGpath)
        if not isExist:  # Create a new directory because it does not exist
            os.makedirs(SHGpath)
            print("The new directory is created!")
            print('--------------------------------------------\n')

        root = Tk()
        root.withdraw()
        folder_selected = filedialog.askdirectory(initialdir="/Users/labaccess/Library/CloudStorage/Box-Box/Jin Lab Shared Folder/SHG-RA data")
        dir_list = os.listdir(folder_selected)
        file_name = dir_list[0]
        print(file_name)
        file_name_1 = dir_list[1]
        for i in range(len(file_name)-1, 0, -1):
            # if file_name[i] == 'p':
            #     file_name = file_name[:i+1]
            if file_name[i] == '_':
                file_name = file_name[:i]
                break
        print()
        for i in range(len(folder_selected) - 1, 0, -1):
            if folder_selected[i] == '/':
                folder_name_pptx = folder_selected[i+1:]
                break

        folder_selected = folder_selected + "/"
        data_sel = 'n'
        Parameter = pd.read_csv(folder_selected + "Experimental_Parameters.txt", header=None, sep=':', engine='c')
        print(Parameter)
        Date = Parameter.iat[0, 1]
        step_size = Parameter.iat[7, 1]
        step_size = int(step_size)
        file_name = str(Parameter.iat[1,1]).replace(" ","")
        print(file_name)
        avg_x = 0
        avg_y = 0
        iteration = 0

        deg_file_org = []
        sig_file_org = []

        # for degree in range(0, 10, step_size):
        degree = 130
        temp_temp = 80
        SHG_Raw = np.loadtxt(folder_selected + file_name + "_{}deg".format(degree) + ".txt", dtype=int, delimiter=',')
        # SHG_Raw = np.loadtxt(folder_selected + file_name + "{}deg".format(degree) + ".txt", dtype=int, delimiter=',')
        # SHG_Raw = np.loadtxt(folder_selected + 'STO_Nb_0_0035_{}k_{}deg'.format(temp_temp, degree) + '.txt', dtype=int, delimiter=',')


        # from scipy.fft import ifftn
        # import matplotlib.pyplot as plt
        # import matplotlib.cm as cm
        # import numpy as np
        #
        # fig, ax = pyplot.subplots()
        # ft = np.fft.ifftshift(SHG_Raw)
        # ft1 = np.fft.fft2(ft)
        # ft2 = np.fft.fftshift(ft1)
        # ft2= np.log(abs(ft2))
        # print(ft2.argmax())
        # print(ft2.argmax()/512)
        # print(ft2.argmax()%512)
        # print(np.max(ft2))
        # ax.imshow(np.log(abs(ft2)))
        # plt.show()

        # SHG_Raw = SHG_Raw[128:384, 128:384]
        fig, ax = pyplot.subplots()
        im = ax.imshow(SHG_Raw, vmin=0, vmax=5000)
        polarization = Parameter.iat[8, 1]
        fig.colorbar(im, ax=ax, label='{} Polarization'.format(polarization))
        exposure_time = str(float(Parameter.iat[9, 1]))
        title = str(Parameter.iat[1, 1]) + '' + str(Parameter.iat[2, 1]) + '' + str(Parameter.iat[3, 1]) \
                + '\n' + str(Parameter.iat[4, 1]) + 'mW Exposure Time ' + exposure_time + 's Averaging ' \
                + str(int(Parameter.iat[11, 1]))
        pyplot.title(title + ' at {} Degree'.format(degree), pad=10, wrap=True)
        # pyplot.show()

        def onclick(event):
            if event.dblclick:
                global cur_x, cur_y, click
                cur_x = event.xdata
                cur_y = event.ydata
                click = event.button
                if click == 1:
                    print('Closing')
                    pyplot.close()
        connection_id = fig.canvas.mpl_connect('button_press_event', onclick)
        pyplot.savefig(folder_selected+"Figure_1.png")
        pyplot.show()

        deg_file = []
        sig_file = []
        center_x = int(cur_x)
        center_y = int(cur_y)
        # center_x = 238
        # center_y = 258
        # center_x = int(input('Enter the center for x axis: '))
        # center_y = int(input('Enter the center for y axis: '))
        region_size = int(input('Enter the box size: '))
        # region_size = 80
        half_region_size = (np.ceil(region_size / 2)).astype(int)
        # min_pos = (postion_max_x + 235 - half_region_size)[0]
        # max_pos = (postion_max_x + 235 + half_region_size)[0]

        SHG_Raw = np.loadtxt(folder_selected + file_name + "_{}deg".format(degree) + ".txt", dtype=int, delimiter=',')

        fig, ax = pyplot.subplots()
            # ax.imshow(SHG_Raw)
        region = SHG_Raw[center_x - half_region_size: center_x + half_region_size,
                             center_y - half_region_size: center_y + half_region_size]
        im = ax.imshow(SHG_Raw, vmin=1000, vmax=5000)

        fig.colorbar(im, ax=ax, label='Interactive colorbar')
        ax.scatter(center_x, center_y, s=30, color='tomato', marker='x')
        rect = patches.Rectangle((center_x - half_region_size, center_y - half_region_size),
                                     region_size, region_size, linewidth=1, edgecolor='r', facecolor='none')
            # Add the patch to the Axes
        ax.add_patch(rect)
        pyplot.show()
        for degree in range(0, 365, step_size):
            deg_file = np.append(deg_file, degree)
            SHG_Raw = np.loadtxt(folder_selected + file_name + "_{}deg".format(degree) + ".txt", dtype=int, delimiter=',')

            region = SHG_Raw[center_x - half_region_size: center_x + half_region_size,
                             center_y - half_region_size: center_y + half_region_size]

            small_sum = sum(map(sum, region))
            large_sum = sum(map(sum, SHG_Raw))
            bkg_pixel = (large_sum - small_sum) / (512 ** 2 - region_size ** 2)
            sig = small_sum - bkg_pixel * region_size ** 2
            sig_file = np.append(sig_file, sig)

        sig_file = sig_file.astype(np.float64)
        # sig_file = sig_file.tolist()
        max_lim = max(sig_file)
        min_lim = min(sig_file)
        deg_file = deg_file * np.pi / 180
        deg_file = deg_file.astype(np.float64)
        # deg_file = deg_file.tolist()
        fig, ax = pyplot.subplots(subplot_kw={'projection': 'polar'})
        ax.plot(deg_file, sig_file, color='red')
        ax.set_ylim(bottom=min_lim, top=max_lim)
        # pyplot.autoscale()
        pyplot.show()
        pyplot.plot(deg_file, sig_file, linewidth=5, color='blue')
        pyplot.show()

        slope = (sig_file[-1] - sig_file[0]) / (deg_file[-1] - deg_file[0])
        const = sig_file[-1] - slope * deg_file[-1]
        # const = sig_file[-1] + slope * deg_file[-1]
        sig_file = sig_file - (slope * deg_file + const)
        sig_file = (sig_file/30)/380000
        csv_file_name = 'Processed_Data.csv'
        comb = pd.DataFrame(list(zip(deg_file, sig_file)))
        rec_data = pd.DataFrame()
        rec_data = pd.concat([rec_data, comb], ignore_index=True, axis=1)
        rec_data.to_csv(folder_selected + csv_file_name, mode='a', index=False, encoding='utf-8-sig', header=None)

        min_sig = min(sig_file)
        sig_file = sig_file - min_sig
        csv_file_name = 'Processed_Data_Min_Shift.csv'
        comb = pd.DataFrame(list(zip(deg_file, sig_file)))
        rec_data = pd.DataFrame()
        rec_data = pd.concat([rec_data, comb], ignore_index=True, axis=1)
        rec_data.to_csv(folder_selected + csv_file_name, mode='w', index=False, encoding='utf-8-sig', header=None)

        pyplot.plot(deg_file, sig_file, linewidth=5, color='blue')
        pyplot.show()

        max_lim = max(sig_file)
        min_lim = min(sig_file)


        fig, ax = pyplot.subplots(subplot_kw={'projection': 'polar'})
        ax.plot(deg_file, sig_file, color='red')
        ax.set_ylim(bottom=min_lim, top=max_lim)
        pyplot.title(title + '{} Polarization'.format(polarization), pad=10, wrap=True)
        pyplot.tight_layout()
        pyplot.savefig(folder_selected+"Figure_2.png")
        pyplot.show()

        fit = True
        if fit == True:
            def shg_sin(params, x, data=None):
                A = params['A']
                a2 = params['a2']
                B = params['B']
                x0 = params['x0']
                model = (A*np.sin(a2-3*(x-x0))+B*np.sin(a2+(x-x0)))**2
                if data is None:
                    return model
                return model - data

            def shg_cos(params, x, data=None):
                A = params['A']
                a2 = params['a2']
                B = params['B']
                model = A*np.cos(a2-3*x)+B*np.cos(a2+3*x)*2
                if data is None:
                    return model
                return model - data

            # Create a Parameters object
            params = lmfit.Parameters()
            params.add('A', value=-0.1)
            params.add('a2', value=-0.1)
            params.add('B', value=11)
            params.add('x0', value=0.1)
            result_sin = lmfit.minimize(shg_sin, params, args=(deg_file,), kws={'data': sig_file})
            # result_cos = lmfit.minimize(shg_cos, params, args=(deg_file,), kws={'data': sig_file})
            # sin_a1 = result_sin.params['a1'].value
            sin_A = result_sin.params['A'].value
            sin_a2 = result_sin.params['a2'].value
            sin_B = result_sin.params['B'].value
            sin_x0 = result_sin.params['x0'].value

            fig, ax = pyplot.subplots(subplot_kw={'projection': 'polar'})
            ax.scatter(deg_file, sig_file, color='black', s=1)
            ax.plot(deg_file, result_sin.residual + sig_file, color='red',linewidth=2)
            ax.set_ylim(bottom=min_lim*1.05, top=max_lim*1.05)
            pyplot.title(title + '{} Polarization'.format(polarization), pad=10, wrap=True)
            pyplot.tight_layout()
            pyplot.savefig(folder_selected + "Fitted_Data.png")
            pyplot.show()

            df = pd.DataFrame()
            df_comb = pd.DataFrame(list(zip([sin_A], [sin_a2], [sin_B], [sin_x0])))
            df = pd.concat([df, df_comb], ignore_index=True, axis=1)
            df.to_csv(folder_selected + 'Fitted_Data.csv', index=False)

        if os.path.exists(folder_selected + 'Results.pptx'):
            prs = Presentation(folder_selected + 'Results.pptx')
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
        else:
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            prs.save(folder_selected + 'Results.pptx')
            prs = Presentation(folder_selected + 'Results.pptx')
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

        SHG_Image = folder_selected + 'Figure_1.png'
        SHG_Signal = folder_selected + 'Figure_2.png'
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        SHG_Image_img = slide.shapes.add_picture(SHG_Image, Inches(0.32), Inches(1.42), Inches(6.39))
        SHG_Signal_img = slide.shapes.add_picture(SHG_Signal, Inches(6.49), Inches(1.42), Inches(6.39))
        text_frame = slide.shapes.add_textbox(Inches(0.18), Inches(0.2), Inches(6.67), Inches(0.4))
        Data_frame = slide.shapes.add_textbox(Inches(11.19), Inches(0.2), Inches(2.04), Inches(0.4))
        text_frame = text_frame.text_frame
        Data_frame = Data_frame.text_frame
        p = text_frame.paragraphs[0]
        d = Data_frame.paragraphs[0]
        run = p.add_run()
        run.text = str(folder_name_pptx)
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(18)

        run_d = d.add_run()
        run_d.text = str(Date)
        font_d = run_d.font
        font_d.name = 'Calibri'
        font_d.size = Pt(18)
        prs.save(folder_selected + 'Results.pptx')


if __name__ == '__main__':
    running = True
    while running:
    # construct the main wi
        try:
            window1 = SHG_Processing()
        except KeyboardInterrupt:
            break
    # endless loop unless quit
